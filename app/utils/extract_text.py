import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import docx
import os
from pathlib import Path
from pptx import Presentation
from openpyxl import load_workbook


def extract_text_from_file(filepath):
    ext = Path(filepath).suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext == ".docx":
        return extract_text_from_docx(filepath)
    elif ext == ".txt":
        return extract_text_from_txt(filepath)
    elif ext == ".pptx":
        return extract_text_from_pptx(filepath)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(filepath)
    elif ext in [".png", ".jpg", ".jpeg"]:
        return extract_text_from_image(filepath)
    else:
        return "Unsupported file type."


def extract_text_from_pdf(filepath):
    doc = fitz.open(filepath)
    text = ""
    for page in doc:
        text += page.get_text()
    return clean_text(text)


def extract_text_from_docx(filepath):
    doc = docx.Document(filepath)
    raw = "\n".join([para.text for para in doc.paragraphs])
    return clean_text(raw)


def extract_text_from_txt(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        return clean_text(f.read())


def extract_text_from_image(filepath):
    img = Image.open(filepath)
    raw = pytesseract.image_to_string(img)
    return clean_text(raw)


def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return clean_text(text)


def extract_text_from_xlsx(filepath):
    wb = load_workbook(filename=filepath, read_only=True)
    text = ""
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
            text += row_text + "\n"
    return clean_text(text)


def clean_text(text):
    return " ".join(text.split()).strip()


def remove_metadata_lines(text):
    lines = text.split("\n")
    filtered = [line for line in lines if not line.lower().startswith(("name:", "index number:"))]
    return " ".join(filtered).strip()